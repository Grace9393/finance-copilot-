"""
image_to_pptx.py
Converts an image file (PNG/JPG/GIF/WebP) into an editable IBM-branded PPTX.

Usage:
  python image_to_pptx.py <image_path> <output_path> [--title "Slide Title"] [--notes "Speaker notes"]

Produces a single slide with:
  - The image filling the left 60% of the slide
  - An editable text frame on the right for annotations
  - Speaker notes containing OCR text (if tesseract is available) or filename
  - IBM Carbon colour theme: slide bg #f4f4f4, accent #0f62fe, text #161616
"""

import sys
import os
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print('{"error": "python-pptx not installed. Run: pip install python-pptx pillow"}', file=sys.stderr)
    sys.exit(1)

# IBM Carbon colours
IBM_BLUE   = RGBColor(0x0F, 0x62, 0xFE)  # #0f62fe
IBM_DARK   = RGBColor(0x16, 0x16, 0x16)  # #161616
IBM_BG     = RGBColor(0xF4, 0xF4, 0xF4)  # #f4f4f4
IBM_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def ocr_text(image_path: str) -> str:
    """Try OCR with pytesseract; return empty string if unavailable."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(image_path)
        return pytesseract.image_to_string(img).strip()
    except Exception:
        return ''


def build_pptx(image_path: str, output_path: str, title: str = '', notes: str = '') -> dict:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Background ──────────────────────────────────────────────────────────────
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = IBM_BG

    # ── Image (left 65% of slide) ────────────────────────────────────────────────
    img_left   = Inches(0.3)
    img_top    = Inches(0.8)
    img_width  = Inches(8.2)
    img_height = Inches(6.2)
    slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)

    # ── Right annotation panel ───────────────────────────────────────────────────
    panel_left  = Inches(8.8)
    panel_top   = Inches(0.8)
    panel_w     = Inches(4.2)
    panel_h     = Inches(6.2)

    txBox = slide.shapes.add_textbox(panel_left, panel_top, panel_w, panel_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Annotation heading
    p_head = tf.paragraphs[0]
    p_head.text = 'Annotations'
    p_head.alignment = PP_ALIGN.LEFT
    run_head = p_head.runs[0]
    run_head.font.size = Pt(14)
    run_head.font.bold = True
    run_head.font.color.rgb = IBM_BLUE

    # Editable placeholder lines
    placeholders = [
        'Key observation 1',
        'Key observation 2',
        'Key observation 3',
        '',
        'Add your notes here…'
    ]
    for text in placeholders:
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.color.rgb = IBM_DARK

    # ── Title bar at top ──────────────────────────────────────────────────────────
    slide_title = title or Path(image_path).stem.replace('_', ' ').replace('-', ' ').title()
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.55))
    tf_title = title_box.text_frame
    tf_title.paragraphs[0].text = slide_title
    if tf_title.paragraphs[0].runs:
        r = tf_title.paragraphs[0].runs[0]
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = IBM_DARK

    # ── Blue accent line under title ──────────────────────────────────────────────
    from pptx.util import Pt as PtLine
    from pptx.oxml.ns import qn
    from lxml import etree
    # Draw a thin blue rectangle as an accent bar
    accent = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.3), Inches(0.67), Inches(12.7), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = IBM_BLUE
    accent.line.fill.background()  # no border

    # ── Speaker notes ────────────────────────────────────────────────────────────
    ocr = ocr_text(image_path)
    notes_text = notes or ocr or f'Image: {Path(image_path).name}'
    slide.notes_slide.notes_text_frame.text = notes_text

    prs.save(output_path)

    return {
        'slide_count': 1,
        'output': output_path,
        'title': slide_title,
        'ocr_chars': len(ocr),
        'warnings': [] if ocr else ['OCR unavailable — install pytesseract for text extraction']
    }


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('image', help='Input image path')
    parser.add_argument('output', help='Output .pptx path')
    parser.add_argument('--title', default='', help='Slide title override')
    parser.add_argument('--notes', default='', help='Speaker notes override')
    args = parser.parse_args()

    import json
    result = build_pptx(args.image, args.output, args.title, args.notes)
    print(json.dumps(result))
